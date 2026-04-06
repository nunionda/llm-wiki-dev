from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_pitch():
    prs = Presentation()

    # Define Colors (Sage/Charcoal/Blue)
    nordic_charcoal = RGBColor(0x1A, 0x1A, 0x1A)
    nordic_white = RGBColor(0xFD, 0xFD, 0xFD)
    nordic_sage = RGBColor(0x8B, 0x9C, 0x8B)
    nordic_blue = RGBColor(0x4A, 0x6F, 0xA5)

    def set_slide_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank
    set_slide_bg(slide, nordic_charcoal)
    
    tx = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    tf = tx.text_frame
    p = tf.add_paragraph()
    p.text = "Agentic Wiki"
    p.font.bold = True
    p.font.size = Pt(72)
    p.font.color.rgb = nordic_white
    p.alignment = PP_ALIGN.LEFT

    p2 = tf.add_paragraph()
    p2.text = "ZERO-DISTANCE KNOWLEDGE ENGINE"
    p2.font.size = Pt(24)
    p2.font.color.rgb = nordic_sage
    p2.alignment = PP_ALIGN.LEFT

    # Slide 2: The Problem
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, nordic_white)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title.text_frame.text = "The Knowledge Decay"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = nordic_charcoal

    body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = body.text_frame
    for point in ["Manual documentation is expensive and dies quickly.", "Human intent (Vibes) is lost in the noise.", "Knowledge siloes are the biggest tax on growth."]:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(24)
        p.font.color.rgb = nordic_charcoal
        p.space_after = Pt(20)

    # Slide 3: The Solution
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, nordic_blue)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    p = title.text_frame.paragraphs[0]
    p.text = "The Autonomous Brain"
    p.font.size = Pt(44)
    p.font.color.rgb = nordic_white

    body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
    tf = body.text_frame
    for step in ["Raw Ingest (Vibe Capturing)", "Deterministic Synthesis", "Knowledge Compounding (Semantic Graph)", "Hybrid Governance"]:
        p = tf.add_paragraph()
        p.text = f"→ {step}"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = nordic_white
        p.space_after = Pt(15)

    # Slide 4: Business Model
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, nordic_white)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    p = title.text_frame.paragraphs[0]
    p.text = "Revenue: Knowledge Computation"
    p.font.size = Pt(44)
    p.font.color.rgb = nordic_charcoal

    # Simple Table Mockup
    rows, cols = 4, 3
    left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(4)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    headers = ["Tier", "Price/Seat", "Value Metric"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    
    data = [
        ["Starter", "Free", "50 Nodes"],
        ["Pro", "$29/mo", "Live Sync + Map"],
        ["Enterprise", "Custom", "On-Premise SLM"]
    ]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            table.cell(r+1, c).text = val

    # Slide 5: Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, nordic_charcoal)
    
    title = slide.shapes.add_textbox(Inches(5), Inches(3), Inches(4), Inches(2))
    p = title.text_frame.paragraphs[0]
    p.text = "JOIN THE ERA OF VIBE CODING."
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = nordic_sage
    p.alignment = PP_ALIGN.RIGHT

    prs.save('Agentic_Wiki_Pitch_Deck.pptx')
    print("Agentic_Wiki_Pitch_Deck.pptx generated successfully.")

if __name__ == "__main__":
    create_pitch()
