from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title, points, notes=""):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title_place = slide.shapes.title
    title_place.text = title
    
    body_place = slide.placeholders[1]
    tf = body_place.text_frame
    tf.word_wrap = True
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        if point.startswith('  -'):
            p.text = point.replace('  -', '').strip()
            p.level = 1
            
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

def generate_manual():
    prs = Presentation()
    
    # 1. Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "gstack Enterprise Agentic SDK"
    subtitle.text = "지능형 자동화와 에이전틱 코딩의 표준 기술 매뉴얼\nv1.0.2 | Technical Report"

    # 2. TOC
    add_slide(prs, "Table of Contents", [
        "01. gstack 개요 및 철학 (Ethos)",
        "02. 아키텍처 및 보안 (Security First)",
        "03. 코어 스킬셋 및 워크플로우 (Lifecycle)",
        "04. 데이터 및 문서 자동화 (Structured Assets)",
        "05. 엔터프라이즈 거버넌스 및 확장 (Audit/MCP)",
        "06. 실전 케이스 스터디 및 모범 사례 (Best Practices)"
    ])

    # 3. Philosophy
    add_slide(prs, "01. gstack Core Philosophy (Ethos)", [
        "Boil the Lake: AI 시대의 생산성을 반영한 100% 완결성 추구",
        "  - 부분 해결이 아닌, 끝까지 파고드는 기술적 집요함",
        "Search Before Building: 지식의 계층화와 중복 제거",
        "  - Standard -> Popular -> First Principles 기반의 의사결정",
        "AI Effort Compression: 불필요한 컨텍스트 낭비 없는 핵심 정보 집중"
    ])

    # 4. Lifecycle
    add_slide(prs, "02. gstack 7-Step Lifecycle", [
        "1. Think: 문제의 본질(Pain Point) 파악",
        "2. Plan: 정교한 implementation_plan.md 수립",
        "3. Build: Bun/TypeScript 기반의 고성능 구현",
        "4. Review: Staff Engineer 레벨의 코드 감사",
        "5. Test: Vitest 및 Playwright 기반의 자동 QA",
        "6. Ship: 원자적 커밋 및 CI/CD 배송",
        "7. Reflect: 운영 학습 및 실무 노하우 기록 (Learn)"
    ])

    # 5. Architecture
    add_slide(prs, "03. Infrastructure: Browser Daemon", [
        "Browser Daemon ($B): 100ms 응답성을 위한 상주형 데몬",
        "  - 세션 쿠키, 탭, 로그인 상태의 영구 유지",
        "병렬 워크스페이스 지원: 10,000~60,000 포트 임의 할당",
        "Localhost Bound: 127.0.0.1 바인딩 및 Bearer Token 인증 강제",
        "Single Binary: Bun Build를 통한 간편한 배포 및 실행"
    ])

    # 6. Security
    add_slide(prs, "04. Enterprise Security: Air-Gapped SLM", [
        "Zero-Trust Architecture: 로컬 전용 추론 및 데이터 소유권 보장",
        "On-Premise 최적화: 외부 대형 모델(LLM) 의존성 최소화",
        "Egress Control: 외부 데이터 유출을 원천 차단하는 보안 가시화",
        "CSO Auditor 모드: OWASP Top 10 + STRIDE 위협 모델링 자동화"
    ])

    # 7. Standards
    add_slide(prs, "05. Engineering Standards & Runtime", [
        "Primary Runtime: Bun (Native SQLite/TypeScript 지원)",
        "Test Framework: Vitest (Unit) / Playwright (E2E)",
        "Strict Type Safety: 'any' 사용 금지 및 기술적 명확성 유지",
        "Atomic Commits: 한 커밋당 하나의 완전한 논리적 생각(Thought) 반영"
    ])

    # 8. Planning (1)
    add_slide(prs, "06. Planning: CEO/Founder Review", [
        "Product Thinking: 10-star 가치를 전달하는 제품 기획",
        "Scope Management: 전략적 확장 및 필수 요소 압축",
        "Business Alignment: 타겟 사용자 페르소나 및 시장 적합성 검토",
        "Skill: plan-ceo-review (Interactive Plan Audit)"
    ])

    # 9. Planning (2)
    add_slide(prs, "07. Planning: Engineering Review", [
        "Architecture First: 데이터 흐름, 다이어그램, 에지 케이스 사전 정의",
        "Performance Design: 지연 시간 최소화 및 확장성 고려",
        "Test Strategy: 유닛 테스트부터 통합 테스트까지의 커버리지 설계",
        "Skill: plan-eng-review (Tech Stack & Security Lock-in)"
    ])

    # 10. Planning (3)
    add_slide(prs, "08. Planning: Design Review", [
        "Visual Excellence: 프리미엄 다크 모드, 글래스모피즘, 동적 애니메이션",
        "Design System: 일관된 컬러 팔레트, 타이포그래피, 간격 규칙 수립",
        "AI Slop 제거: 제너릭한 UI를 지양하고 창의적/미학적 인터페이스 구현",
        "Skill: plan-design-review (High-Fidelity Mockup Audit)"
    ])

    # 11. Building (1)
    add_slide(prs, "09. Building: Frontend Finalization", [
        "design-html: Mockup을 실제 프로덕션 레벨의 HTML/CSS로 전환",
        "Zero Dependency: 30KB 미만의 오버헤드를 가진 경량/고성능 코드",
        "Smart API Routing: Pretext 패턴 기반의 최적 인터페이스 라우팅",
        "Responsive Layout: 모바일-데스크탑 하이브리드 지원"
    ])

    # 12. Building (2)
    add_slide(prs, "10. Building: AI Logic Integration", [
        "Claude API & Agent SDK: 지능형 자율 행동을 위한 SDK 연동",
        "MCP Builder: 외부 서비스 연동을 위한 Model Context Protocol 구축",
        "Agentic Wiki Sync: 위키 지식을 기반으로 한 RAG/컴파운딩 로직",
        "Skill: claude-api, mcp-builder, web-artifacts-builder"
    ])

    # 13. QA (1)
    add_slide(prs, "11. Quality Assurance: Systemic Testing", [
        "QA Tiers: Quick(Critical), Standard(Medium), Exhaustive(Full Audit)",
        "Bug Repro: Playwright 기반의 버그 재현 및 자동 픽스 루프",
        "Fix Evidence: 수정 전/후 스크린샷 및 상태 비교 리포트 제공",
        "Skill: qa (Test-Fix-Verify Loop)"
    ])

    # 14. QA (2)
    add_slide(prs, "12. Browser & Site Dogfooding", [
        "Headless Browse: JS 실행, 인증(Cookie), 폼 처리 지원",
        "Gstack Browser Extension: AI 전용 사이드바가 포함된 실시간 모니터링",
        "Accessibility Audit: 장애물 없는 인터페이스 및 웹 표준 준수 확인",
        "Skill: browse, setup-browser-cookies, open-gstack-browser"
    ])

    # 15. Deployment
    add_slide(prs, "13. Deployment: The Ship Workflow", [
        "Atomic Delivery: 테스트-리뷰-배포가 하나로 묶인 원스톱 배송",
        "Base Branch Sync: 최신 코드 동기화 및 충돌 방지 로직",
        "Version Control: 자동 버전 범프 및 CHANGELOG 업데이트",
        "Skill: ship (PR Creation & Release Automation)"
    ])

    # 16. Post-Ship
    add_slide(prs, "14. Monitoring: Performance & Canary", [
        "Benchmark: Core Web Vitals 및 페이지 로딩 성능 추적",
        "Canary: 배포 후 실시간 콘솔 오류 및 상태 이상 감지",
        "Health Dashboard: 코드 품질 및 시스템 안정성 점수화(0-10)",
        "Skill: benchmark, canary, health"
    ])

    # 17. Structured Data (1)
    add_slide(prs, "15. Data Assets: Excel & Docs Automation", [
        "XLSX Automation: 대용량 데이터 정제, 수식 계산, 차트 자동 생성",
        "DOCX/PDF Reporting: 고품질 기술 문서 및 가이드북 자동 레이아웃",
        "Data Analyst Expert: Causal Inference 기반의 비즈니스 인사이트 도출",
        "Skill: xlsx, docx, pdf, senior-data-analyst"
    ])

    # 18. Structured Data (2)
    add_slide(prs, "16. Presentation: Automated PPTX Decks", [
        "PPTX Skill: 피치 데크, 보고용 장표, 매뉴얼의 자동 생성",
        "Slide Layout Intelligence: 텍스트 밀도와 시각 요소의 최적 배치",
        "Theme Factory: 10가지 프리셋 테마 기반의 자동 스타일링",
        "Skill: pptx, theme-factory"
    ])

    # 19. Knowledge Sync (1)
    add_slide(prs, "17. Knowledge: NotebookLM Integration", [
        "Source Grounding: 사용자 보유 지식을 기반으로 한 환각(Hallucination) 제로",
        "Research Partner: 수만 단어의 문서를 읽고 전문적 대화 및 가이드 생성",
        "Context Management: 세션 기반의 깊이 있는 연구 파트너십",
        "Skill: mcp-notebooklm (ask_question, list_notebooks)"
    ])

    # 20. Knowledge Sync (2)
    add_slide(prs, "18. Knowledge: Agentic Wiki Engine", [
        "Vibe Ingest: 파편화된 메모의 즉각적인 지식화 및 카테고리 분류",
        "Compounding Intelligence: 시간이 흐를수록 정제되는 지식 베이스",
        "Semantic Matrix: 지식 간의 관계를 분석하여 전략 사각지대 탐지",
        "Action: actionable-refactor, gen-manifest"
    ])

    # 21. Operational Excellence
    add_slide(prs, "19. Governance: CSO Audit Mode", [
        "Security Archaeology: 하드코딩된 Secret 및 취약한 라이브러리 탐지",
        "Chain of Custody: 인프라 변경 이력 및 보안 무결성 추적",
        "Guard & Freeze Mode: 작업 범위를 제한하여 파괴적 변경 방지",
        "Skill: cso, careful, guard, freeze"
    ])

    # 22. DX & Collaboration
    add_slide(prs, "20. Developer Experience (DX) Focus", [
        "TTHW Optimization: 'Hello World' 도달 시간 최소화 설계",
        "Internal Comms: 팀 보고서, 리더십 업데이트용 통신 포맷 지원",
        "Office Hours: 디자인 씽킹 기반의 브레인스토밍 및 디자인 문서화",
        "Skill: devex-review, internal-comms, office-hours"
    ])

    # 23. Customization
    add_slide(prs, "21. Extensibility: Creating New Skills", [
        "Skill Creator: gstack 생태계를 확장하기 위한 고유 스킬 개발",
        "Template Based: 일관된 인터페이스를 가진 스킬 구현 가이드",
        "MCP SDK: Node.js/Python 기반의 모델 컨텍스트 프로토콜 연동",
        "Skill: skill-creator, mcp-builder"
    ])

    # 24. Performance
    add_slide(prs, "22. Optimization & Health Management", [
        "Retro: 주간 커밋 분석 및 팀 기여도 가시화",
        "Learn: 시행착오 및 운영 노하우의 영구 보존",
        "Upgrade: 최신 gstack 툴링으로의 자동 마이그레이션",
        "Skill: retro, learn, upgrade, health"
    ])

    # 25. Best Practices (1)
    add_slide(prs, "23. The Iron Law of Debugging", [
        "Hypothesis Driven: 가정 수립 -> 데이터 트레이스 -> 근본 원인(RCA) 분석",
        "No YOLO Fix: 원인을 알지 못하는 상태에서의 수정(YOLO) 절대 금지",
        "Investigate Mode: 4단계 정밀 디버깅 루프 수행",
        "Skill: investigate (/analyze, /hypothesize)"
    ])

    # 26. Best Practices (2)
    add_slide(prs, "24. Professional Builder Voice", [
        "Technical Clarity: 명확한 데이터 기반의 소통",
        "No AI Slop: 의미 없는 사과나 감탄사 배제",
        "Decision Artifacts: 의사결정 과정을 문서화(Plan/Task/Walkthrough)하여 투명성 확보"
    ])

    # 27. Case Study (1)
    add_slide(prs, "25. Case Study: Rapid Web Prototyping", [
        "기획: Office Hours를 통한 아이디어 구체화",
        "설계: CEO/Eng Review를 통해 10-star 가치 및 아키텍처 확정",
        "구현: Frontend Design & HTML Finalization",
        "QA/Ship: 10분 내에 프로덕션 배포 완료"
    ])

    # 28. Case Study (2)
    add_slide(prs, "26. Case Study: Enterprise Security Infrastructure", [
        "감사: CSO Mode를 통해 온프레미스 인프라 취약점 스캔",
        "구축: Air-Gapped 보안 배지 및 신뢰 대시보드 적용",
        "검증: Canary & Benchmark를 통한 실시간 안정성 입증"
    ])

    # 29. Roadmap
    add_slide(prs, "27. The Future: Agentic Operating System", [
        "Self-Healing Code: 스스로 버그를 고치는 자가 치유 시스템",
        "Distributed Brain: 여러 에이전트가 협업하는 분산형 지능 구조",
        "Real-time Knowledge Web: 전 세계 지식이 실시간으로 동기화되는 위키"
    ])

    # 30. Conclusion
    add_slide(prs, "28. Conclusion: Beyond AI Tools to AI Agents", [
        "gstack은 단순한 도구가 아닌 사용자의 생각을 실행으로 옮기는 '뇌'입니다.",
        "Boil the Lake 철학으로 완성도 높은 미래를 함께 만듭니다.",
        "Contact: gstack support team | v1.0.2 FINAL"
    ])

    # Save
    target_file = "gstack_Manual_Guide.pptx"
    prs.save(target_file)
    print(f"Generated: {target_file}")

if __name__ == "__main__":
    generate_manual()
